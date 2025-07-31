import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd

def parse_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def parse_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def parse_file(file_path):
    if file_path.endswith('.pdf'):
        return parse_pdf(file_path)
    elif file_path.endswith('.docx'):
        return parse_docx(file_path)
    elif file_path.endswith('.pptx'):
        return parse_pptx(file_path)
    elif file_path.endswith('.csv'):
        return parse_csv(file_path)
    elif file_path.endswith('.txt') or file_path.endswith('.md'):
        return parse_txt(file_path)
    else:
        return "Unsupported file format"
